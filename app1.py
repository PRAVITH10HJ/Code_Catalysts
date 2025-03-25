import os
import cv2
import pytesseract
import nltk
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches, Pt
from dotenv import load_dotenv
import textwrap

# ✅ Load environment variables
load_dotenv()

# ✅ Initialize Flask app
app = Flask(__name__)
CORS(app)

UPLOAD_FOLDER = "uploads"
OUTPUT_FOLDER = "output"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

# ✅ Ensure NLTK tokenizer is available
try:
    nltk.data.find("tokenizers/punkt")
except LookupError:
    nltk.download("punkt")

# 🔹 Function to Extract Text Using Tesseract OCR
def process_image_tesseract(image_path):
    """Extract text from an image using Tesseract OCR"""
    img = cv2.imread(image_path)
    if img is None:
        return "Error: Image could not be loaded."

    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    text = pytesseract.image_to_string(gray)
    return text.strip()

# 🔹 Function to Create PowerPoint Slides with Proper Formatting
def create_pptx(text, output_path):
    """Generate a PowerPoint file from extracted text with proper formatting."""
    prs = Presentation()
    sentences = nltk.sent_tokenize(text)

    if not sentences:
        return

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = "Generated Presentation"
    subtitle.text = "Converted from Image Text"

    max_chars_per_slide = 400  # ✅ Limit content per slide
    wrapped_sentences = []

    for sentence in sentences:
        wrapped_sentences.extend(textwrap.wrap(sentence, width=80))  # ✅ Wrap text to avoid overflowing

    slide_text = ""
    for line in wrapped_sentences:
        if len(slide_text) + len(line) > max_chars_per_slide:
            # ✅ Add a new slide when the text exceeds max limit
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content layout
            text_box = slide.shapes.placeholders[1]  # Content box
            text_frame = text_box.text_frame
            text_frame.text = slide_text
            text_frame.paragraphs[0].font.size = Pt(24)  # Adjust font size
            slide_text = ""

        slide_text += line + "\n"

    # ✅ Add remaining text to the last slide
    if slide_text:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        text_box = slide.shapes.placeholders[1]
        text_frame = text_box.text_frame
        text_frame.text = slide_text
        text_frame.paragraphs[0].font.size = Pt(24)

    prs.save(output_path)

# ✅ Define Flask route
@app.route("/upload", methods=["POST"])
def upload_file():
    """Handle file upload, process image with OCR, and generate PowerPoint."""
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "Empty filename"}), 400

    file_path = os.path.join(UPLOAD_FOLDER, file.filename)
    file.save(file_path)

    extracted_text = process_image_tesseract(file_path)
    if "Error" in extracted_text:
        return jsonify({"error": extracted_text}), 500

    output_pptx = os.path.join(OUTPUT_FOLDER, "output_presentation.pptx")
    create_pptx(extracted_text, output_pptx)

    return send_file(output_pptx, as_attachment=True, download_name="generated_presentation.pptx")

# ✅ Run the Flask app
if __name__ == "__main__":
    app.run(debug=True)
